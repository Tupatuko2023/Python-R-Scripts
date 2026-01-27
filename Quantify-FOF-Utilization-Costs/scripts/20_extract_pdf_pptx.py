#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable, List, Optional, Sequence

from path_resolver import get_data_root, get_paper02_dir

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DERIVED_DIR = PROJECT_ROOT / "docs" / "derived_text"
DEFAULT_JSONL = DERIVED_DIR / "chunks.jsonl"
LOG_PATH = DERIVED_DIR / "extract_log.json"


@dataclass
class Chunk:
    source_file: str
    page: Optional[int]
    chunk_type: str  # text | table | slide_text | slide_table
    content_md: str
    created_at: str


def utc_now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()


def write_log(status: str, message: str, extra: Optional[dict[str, Any]] = None) -> None:
    DERIVED_DIR.mkdir(parents=True, exist_ok=True)
    payload: dict[str, Any] = {"timestamp": utc_now_iso(), "status": status, "message": message}
    if extra:
        payload.update(extra)
    LOG_PATH.write_text(json.dumps(payload, ensure_ascii=True, indent=2), encoding="utf-8")


def _safety_check_no_identifier(text: str) -> None:
    low = text.lower()
    if "id," in low or " id " in low:
        raise SystemExit("Safety check failed: identifier-like token detected in derived output.")


def table_to_markdown(table: Sequence[Sequence[Any]]) -> str:
    rows: List[List[str]] = []
    for r in table:
        rows.append([("" if c is None else str(c)).strip() for c in r])
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    for r in rows:
        if len(r) < width:
            r.extend([""] * (width - len(r)))

    header = rows[0]
    sep = ["---"] * width
    body = rows[1:] if len(rows) > 1 else []

    def fmt_row(r: List[str]) -> str:
        return "| " + " | ".join(r) + " |"

    out = [fmt_row(header), fmt_row(sep)]
    out.extend(fmt_row(r) for r in body)
    return "\n".join(out)


def write_jsonl(path: Path, chunks: Iterable[Chunk]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        for ch in chunks:
            f.write(json.dumps(asdict(ch), ensure_ascii=True) + "\n")


def extract_pdf_chunks(pdf_path: Path) -> List[Chunk]:
    created = utc_now_iso()

    try:
        import pdfplumber  # type: ignore
    except Exception:
        pdfplumber = None  # type: ignore

    if pdfplumber is not None:
        out: List[Chunk] = []
        with pdfplumber.open(str(pdf_path)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                txt = (page.extract_text() or "").strip()
                if txt:
                    _safety_check_no_identifier(txt)
                    out.append(Chunk(str(pdf_path), i, "text", txt, created))
                try:
                    tables = page.extract_tables() or []
                except Exception:
                    tables = []
                for t in tables:
                    md = table_to_markdown(t)
                    if md.strip():
                        _safety_check_no_identifier(md)
                        out.append(Chunk(str(pdf_path), i, "table", md, created))
        return out

    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        PdfReader = None  # type: ignore

    if PdfReader is not None:
        out2: List[Chunk] = []
        r = PdfReader(str(pdf_path))
        for idx, page in enumerate(r.pages, start=1):
            try:
                txt = (page.extract_text() or "").strip()
            except Exception:
                txt = ""
            if txt:
                _safety_check_no_identifier(txt)
                out2.append(Chunk(str(pdf_path), idx, "text", txt, created))
        return out2

    raise RuntimeError("No PDF parser available (install pdfplumber or pypdf).")


def extract_pptx_chunks(pptx_path: Path) -> List[Chunk]:
    created = utc_now_iso()
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise RuntimeError("python-pptx not available (install python-pptx).") from e

    prs = Presentation(str(pptx_path))
    out: List[Chunk] = []
    for sidx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                t = str(shape.text).strip()
                if t:
                    texts.append(t)
        if texts:
            joined = "\n".join(texts)
            _safety_check_no_identifier(joined)
            out.append(Chunk(str(pptx_path), sidx, "slide_text", joined, created))

        for shape in slide.shapes:
            if not hasattr(shape, "has_table") or not shape.has_table:
                continue
            tbl = shape.table
            grid: List[List[str]] = []
            for r in tbl.rows:
                row: List[str] = []
                for c in r.cells:
                    row.append((c.text or "").strip())
                grid.append(row)
            md = table_to_markdown(grid)
            if md.strip():
                _safety_check_no_identifier(md)
                out.append(Chunk(str(pptx_path), sidx, "slide_table", md, created))
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description="Layout-aware PDF/PPTX extraction into JSONL chunks (safe-by-default).")
    ap.add_argument("--scan", default="paper_02", help="Scan target (default: paper_02).")
    ap.add_argument("--out", default=str(DEFAULT_JSONL), help="Output JSONL path (default: docs/derived_text/chunks.jsonl)")
    ap.add_argument("--limit-files", type=int, default=50, help="Max number of files to process per run.")
    ap.add_argument("--dry-run", action="store_true", help="Do not write outputs, only log actions.")
    args = ap.parse_args()

    data_root = get_data_root(require=False)
    if not data_root or not data_root.exists():
        write_log("SKIPPED", "DATA_ROOT not set; nothing to extract.")
        print("DATA_ROOT not set; nothing to extract.")
        return 0

    if args.scan != "paper_02":
        write_log("ERROR", f"Unknown scan target: {args.scan}")
        print(f"Unknown scan target: {args.scan}")
        return 2

    base = get_paper02_dir(data_root)
    if not base.exists():
        write_log("SKIPPED", f"Scan base missing: {base}")
        print(f"Scan base missing: {base}")
        return 0

    pdfs = sorted(base.glob("**/*.pdf"))
    pptxs = sorted(base.glob("**/*.pptx"))
    files = (pdfs + pptxs)[: max(0, int(args.limit_files))]

    if not files:
        write_log("SKIPPED", f"No PDF/PPTX files found under: {base}")
        print(f"No PDF/PPTX files found under: {base}")
        return 0

    all_chunks: List[Chunk] = []
    errors: List[str] = []
    for fp in files:
        try:
            if fp.suffix.lower() == ".pdf":
                all_chunks.extend(extract_pdf_chunks(fp))
            elif fp.suffix.lower() == ".pptx":
                all_chunks.extend(extract_pptx_chunks(fp))
        except Exception as e:
            errors.append(f"{fp}: {e}")

    if args.dry_run:
        write_log("DRY_RUN", f"Would write {len(all_chunks)} chunks.", {"errors": errors[:20]})
        print(f"DRY_RUN: would write {len(all_chunks)} chunks")
        return 0

    out_path = Path(args.out)
    write_jsonl(out_path, all_chunks)
    write_log("OK", f"Wrote {len(all_chunks)} chunks.", {"out": str(out_path), "errors": errors[:50]})
    print(f"Wrote {len(all_chunks)} chunks to {out_path}")
    if errors:
        print("Some files failed (see extract_log.json).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
